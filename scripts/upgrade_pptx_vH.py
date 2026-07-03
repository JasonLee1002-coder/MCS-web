
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

path_in = r'G:\我的雲端硬碟\2025_銓幻元_MCS相關資料\07_行銷素材\公司簡報\東方美陳董簡報_20260625\東方美陳董簡報_20260625_20260625_vG.pptx'
path_out = r'G:\我的雲端硬碟\2025_銓幻元_MCS相關資料\07_行銷素材\公司簡報\東方美陳董簡報_20260625\東方美陳董簡報_20260625_20260625_vH.pptx'

prs = Presentation(path_in)

FILL_MAP = {
    '1A1500': '0D1B2A',
    '1A0808': '0D1B2A',
    '8B0000': '0D1B2A',
    '12102A': '0D1B2A',
    '1A2A10': 'EBF5E8',
    '1A2010': 'EBF5E8',
    '20301A': 'EBF5E8',
    'EEF5E8': 'F0FBF0',
    'F5F5E8': 'FAFDF5',
}

TEXT_MAP = {
    'FFFFFF': '0D1B2A',
    '374151': '475569',
    '4ADE80': '16A34A',
    '00C6AD': '00A896',
    'EF4444': 'DC2626',
    'F59E0B': 'D97706',
}

DARK_BG_FILLS = {'0D1B2A', 'C9A84C', '00C6AD', '00A896', '16A34A', 'DC2626', 'D97706'}

def scale_pt(pt):
    if pt < 11:
        new = pt * 1.45
    elif pt < 14:
        new = pt * 1.25
    elif pt < 18:
        new = pt * 1.15
    elif pt < 24:
        new = pt * 1.10
    else:
        new = pt
    return round(new * 2) / 2

def get_fill_hex(shape):
    try:
        if shape.fill.type == 1:
            return str(shape.fill.fore_color.rgb).upper()
    except:
        pass
    return None

stats = {'fills': 0, 'text_colors': 0, 'font_sizes': 0, 'bold': 0, 'lines': 0}

for slide_idx, slide in enumerate(prs.slides):
    # STEP 5: Set white background
    try:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    except:
        pass

    for shape in slide.shapes:
        # STEP 1: Fix fills
        fill_hex = None
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type == 1:
                    c = str(shape.fill.fore_color.rgb).upper()
                    fill_hex = c
                    if c in FILL_MAP:
                        new_c = FILL_MAP[c]
                        shape.fill.fore_color.rgb = RGBColor.from_string(new_c)
                        fill_hex = new_c
                        stats['fills'] += 1
            except:
                pass

        # STEP 6: Fix line/connector colors
        for st in (5, 9):  # FREEFORM=5, LINE=9
            if shape.shape_type == st:
                try:
                    lc = str(shape.line.color.rgb).upper()
                    if lc in ('FFFFFF', 'F5F0E8', 'E8E0D8', 'B8B0A0', 'A09890', '374151'):
                        shape.line.color.rgb = RGBColor.from_string('C9A84C')
                        stats['lines'] += 1
                except:
                    pass

        # STEP 2 & 3 & 4: Text colors, sizes, bold
        on_dark_bg = fill_hex in DARK_BG_FILLS if fill_hex else False

        if hasattr(shape, 'text_frame'):
            all_runs = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt_val = run.font.size / 12700
                        all_runs.append((para, run, pt_val))

            for para, run, pt in all_runs:
                # STEP 2: Color mapping
                try:
                    if run.font.color and run.font.color.type is not None:
                        c = str(run.font.color.rgb).upper()
                        if on_dark_bg:
                            if c in ('0D1B2A', '475569', '374151'):
                                run.font.color.rgb = RGBColor.from_string('FFFFFF')
                                stats['text_colors'] += 1
                        else:
                            if c in TEXT_MAP:
                                run.font.color.rgb = RGBColor.from_string(TEXT_MAP[c])
                                stats['text_colors'] += 1
                except Exception:
                    pass

                # STEP 3: Font size scaling
                new_pt = scale_pt(pt)
                if abs(new_pt - pt) > 0.1:
                    run.font.size = Pt(new_pt)
                    stats['font_sizes'] += 1

                # STEP 4: Bold for title text (>= 22pt)
                if new_pt >= 22 and run.font.bold is None:
                    run.font.bold = True
                    stats['bold'] += 1

# STEP 7: Fix P15 specifically (slide index 14)
if len(prs.slides) > 14:
    slide_p15 = prs.slides[14]
    for shape in slide_p15.shapes:
        fill_hex = get_fill_hex(shape)
        on_dark = fill_hex in DARK_BG_FILLS if fill_hex else False
        if hasattr(shape, 'text_frame') and not on_dark:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.type is not None:
                            c = str(run.font.color.rgb).upper()
                            if c == 'FFFFFF':
                                run.font.color.rgb = RGBColor.from_string('0D1B2A')
                                stats['text_colors'] += 1
                    except:
                        pass

prs.save(path_out)
print('SAVED:', path_out)
size_mb = os.path.getsize(path_out) / (1024*1024)
print(f'File size: {size_mb:.2f} MB')
print()
print('=== Stats ===')
for k, v in stats.items():
    print(f'  {k}: {v} changes')

# Verification
print()
print('=== Verification ===')
prs2 = Presentation(path_out)

def check_slide_text(slide, label):
    issues = []
    for shape in slide.shapes:
        fh = get_fill_hex(shape)
        on_dark = fh in DARK_BG_FILLS if fh else False
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.type is not None:
                            c = str(run.font.color.rgb).upper()
                            if c == 'FFFFFF' and not on_dark:
                                issues.append(f'  FFFFFF in light shape (fill={fh}): "{run.text[:30]}"')
                    except:
                        pass
    if issues:
        print(f'{label}: {len(issues)} rogue FFFFFF found')
        for i in issues[:3]:
            print(i)
    else:
        print(f'{label}: OK - no rogue FFFFFF text on light bg')

slides = prs2.slides
for idx, label in [(0, 'Slide 1'), (9, 'Slide 10'), (11, 'Slide 12'), (14, 'Slide 15 (P15)')]:
    if len(slides) > idx:
        check_slide_text(slides[idx], label)

# Check bad fills gone
print()
print('=== Fill Color Verification ===')
bad_fills = {'1A1500', '1A0808', '8B0000', '12102A', '1A2A10', '1A2010', '20301A'}
found_bad = []
for si, slide in enumerate(prs2.slides):
    for shape in slide.shapes:
        fh = get_fill_hex(shape)
        if fh and fh.upper() in bad_fills:
            found_bad.append(f'Slide {si+1}: {fh}')
if found_bad:
    print(f'BAD fills still present: {found_bad[:5]}')
else:
    print('All rogue dark fills successfully converted!')

print()
print('Done.')
