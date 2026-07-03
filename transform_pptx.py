#!/usr/bin/env python3
"""
Transform 東方美陳董簡報 from dark navy theme to bright/light theme
and rebuild slide 15 (P15) with new financial story layout.
"""

import os
import copy
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

INPUT_PATH = "G:/我的雲端硬碟/2025_銓幻元_MCS相關資料/07_行銷素材/公司簡報/東方美陳董簡報_20260625/東方美陳董簡報_20260625_20260625_vF.pptx"
OUTPUT_PATH = "G:/我的雲端硬碟/2025_銓幻元_MCS相關資料/07_行銷素材/公司簡報/東方美陳董簡報_20260625/東方美陳董簡報_20260625_20260625_vG.pptx"
IMG1_PATH = "C:/Users/JasonLee/claude_code_projects/CMO/tmp_slide_img1.png"
IMG2_PATH = "C:/Users/JasonLee/claude_code_projects/CMO/tmp_slide_img2.png"
IMG3_PATH = "C:/Users/JasonLee/claude_code_projects/CMO/tmp_slide_img3.png"

# Dark → Light background color map
BG_COLOR_MAP = {
    "152338": "FFFFFF",
    "0A1628": "FFFFFF",
    "1A2E4A": "F0F4F8",
    "0F1E35": "F0F4F8",
    "1E3256": "EEF2F7",
    "162440": "EEF2F7",
    "2A4A6F": "E0EAF4",
    "1E3A5F": "E0EAF4",
    "104040": "E0F4F0",
    "0D2A2A": "E0F4F0",
    "133020": "EBF4E8",
    "0D1F18": "EBF4E8",
    "20301A": "EEF5E8",
    "1A2010": "EEF5E8",
    "28281A": "F5F5E8",
    "1A1A10": "F5F5E8",
    "223318": "EEF5E8",
    "1A2A10": "EEF5E8",
    "261C00": "FFF8E8",
    "1A1200": "FFF8E8",
}

# Light → Dark text color map
TEXT_COLOR_MAP = {
    "FFFFFF": "0D1B2A",
    "E8E0D8": "374151",
    "D0C8C0": "6B7280",
    "F5F0E8": "0D1B2A",
    "B8B0A0": "374151",
    "A09890": "6B7280",
}

# Colors to keep unchanged
KEEP_COLORS = {"C9A84C", "00C6AD", "4ADE80", "EF4444", "F59E0B", "A78BFA"}


def normalize_hex(hex_str):
    """Normalize hex string to uppercase 6 chars."""
    return hex_str.upper().lstrip('#')


def set_slide_background_white(slide):
    """Set slide background to white."""
    try:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    except Exception as e:
        print(f"  Warning: Could not set background: {e}")


def delete_fullslide_background_images(slide):
    """Delete full-slide background PICTURE shapes."""
    shapes_to_delete = []
    try:
        sp_tree = slide.shapes._spTree
        for child in list(sp_tree):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'pic':
                # Get position/size from spPr/xfrm
                try:
                    xfrm = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                    if xfrm is not None:
                        ext = xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                        if ext is not None:
                            cx = int(ext.get('cx', 0))
                            cy = int(ext.get('cy', 0))
                            if cx > 10000000 and cy > 5000000:
                                shapes_to_delete.append(child)
                                print(f"    Marking full-slide image for deletion: cx={cx}, cy={cy}")
                except Exception:
                    pass
        for s in shapes_to_delete:
            s.getparent().remove(s)
    except Exception as e:
        print(f"  Warning: Error deleting background images: {e}")


def transform_fill_color(shape_elem):
    """Transform dark fill colors to light ones in XML element."""
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Find all solid fill elements
    for solidFill in shape_elem.findall(f'.//{{{NS_A}}}solidFill'):
        srgbClr = solidFill.find(f'{{{NS_A}}}srgbClr')
        if srgbClr is not None:
            val = srgbClr.get('val', '').upper()
            if val in BG_COLOR_MAP:
                new_val = BG_COLOR_MAP[val]
                srgbClr.set('val', new_val)


def transform_text_colors(shape_elem):
    """Transform light text colors to dark ones in XML element."""
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Find all solid fills in run properties (text colors)
    for rPr in shape_elem.findall(f'.//{{{NS_A}}}rPr'):
        solidFill = rPr.find(f'{{{NS_A}}}solidFill')
        if solidFill is not None:
            srgbClr = solidFill.find(f'{{{NS_A}}}srgbClr')
            if srgbClr is not None:
                val = srgbClr.get('val', '').upper()
                if val in TEXT_COLOR_MAP and val not in KEEP_COLORS:
                    new_val = TEXT_COLOR_MAP[val]
                    srgbClr.set('val', new_val)


def process_all_slides(prs):
    """Apply global bright theme to all slides."""
    total = len(prs.slides)
    for slide_idx, slide in enumerate(prs.slides):
        print(f"  Processing slide {slide_idx + 1}/{total}...")

        # A) Set white background
        set_slide_background_white(slide)

        # B) Delete full-slide background images
        delete_fullslide_background_images(slide)

        # C & D) Transform colors in all XML elements
        sp_tree = slide.shapes._spTree
        for child in sp_tree:
            transform_fill_color(child)
            transform_text_colors(child)


# ============================================================
# Helper functions for building shapes
# ============================================================

def add_rect(slide, left, top, width, height, fill_hex, border_hex=None, border_width_pt=0):
    """Add a solid-fill rectangle to slide."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1=RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if border_hex:
        shape.line.color.rgb = RGBColor.from_string(border_hex)
        shape.line.width = Pt(border_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, color_hex, size_pt, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    """Add a textbox with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color_hex)
    return txBox


def add_textbox_multiline(slide, left, top, width, height, lines):
    """Add a textbox with multiple lines. lines = [(text, color_hex, size_pt, bold, align), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color_hex, size_pt, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color_hex)
    return txBox


def add_image(slide, img_path, left, top, width, height):
    """Add image to slide at specified position."""
    if os.path.exists(img_path):
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        return pic
    else:
        print(f"  Warning: Image not found: {img_path}")
        return None


def add_two_col_row(slide, left, top, width, height, left_text, right_text,
                    left_color, right_color, left_size, right_size,
                    left_bold=False, right_bold=False, fill_hex=None):
    """Add a row with left text and right text in two columns."""
    if fill_hex:
        add_rect(slide, left, top, width, height, fill_hex)

    col_left_w = int(width * 0.65)
    col_right_w = width - col_left_w

    add_textbox(slide, left + 50000, top, col_left_w - 50000, height,
                left_text, left_color, left_size, left_bold)
    add_textbox(slide, left + col_left_w, top, col_right_w - 50000, height,
                right_text, right_color, right_size, right_bold, align=PP_ALIGN.RIGHT)


def rebuild_p15(slide, img1_path, img2_path, img3_path):
    """Completely rebuild slide 15 content."""
    print("  Rebuilding slide 15...")

    # ---- Step 1: Delete unwanted shapes ----
    # Keep only these by name: Rectangle 3, Rectangle 4, Rectangle 5, Rectangle 6,
    # TextBox 7, TextBox 8, Rectangle 9, TextBox 57 (page number)
    KEEP_NAMES = {'Rectangle 3', 'Rectangle 4', 'Rectangle 5', 'Rectangle 6',
                  'TextBox 7', 'TextBox 8', 'Rectangle 9', 'TextBox 57'}

    sp_tree = slide.shapes._spTree
    shapes_to_remove = []

    for child in list(sp_tree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'grpSp'):
            # Get name
            cNvPr = child.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if cNvPr is None:
                # try drawingML namespace
                for ns in ['http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing',
                           'http://schemas.openxmlformats.org/presentationml/2006/main']:
                    cNvPr = child.find(f'{{{ns}}}cNvPr')
                    if cNvPr is not None:
                        break

            # Also check with pptx ns
            if cNvPr is None:
                # Direct children name search
                for desc in child.iter():
                    local = desc.tag.split('}')[-1] if '}' in desc.tag else desc.tag
                    if local == 'cNvPr':
                        cNvPr = desc
                        break

            name = cNvPr.get('name', '') if cNvPr is not None else ''

            if name not in KEEP_NAMES and name != '':
                shapes_to_remove.append(child)
            elif name == '':
                # Unknown element, check if it's an image (pic tag)
                if tag == 'pic':
                    shapes_to_remove.append(child)

    print(f"    Removing {len(shapes_to_remove)} shapes from P15")
    for s in shapes_to_remove:
        if s.getparent() is not None:
            s.getparent().remove(s)

    # ---- Step 2: Update title textbox ----
    # Find TextBox 8 and update title
    for child in sp_tree:
        cNvPr = None
        for desc in child.iter():
            local = desc.tag.split('}')[-1] if '}' in desc.tag else desc.tag
            if local == 'cNvPr':
                cNvPr = desc
                break
        if cNvPr is not None and cNvPr.get('name') == 'TextBox 8':
            # Update text
            NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            # Clear existing text and set new
            txBody = child.find(f'.//{{{NS_A}}}txBody')
            if txBody is not None:
                for p_elem in txBody.findall(f'{{{NS_A}}}p'):
                    for r_elem in p_elem.findall(f'{{{NS_A}}}r'):
                        t = r_elem.find(f'{{{NS_A}}}t')
                        if t is not None:
                            t.text = "24H 自助升級　房租不動，收入 +46%"
                            # Update color to dark
                            rPr = r_elem.find(f'{{{NS_A}}}rPr')
                            if rPr is not None:
                                solidFill = rPr.find(f'{{{NS_A}}}solidFill')
                                if solidFill is not None:
                                    srgbClr = solidFill.find(f'{{{NS_A}}}srgbClr')
                                    if srgbClr is not None:
                                        srgbClr.set('val', '0D1B2A')
                        break
                    break
            break

    # ---- Step 3: Add new content shapes ----

    # 1. Section divider line (gold horizontal bar)
    add_rect(slide, 609600, 1050000, 10972800, 12700, 'C9A84C')

    # 2. LEFT IMAGE BLOCK
    # Image 1: 傳統早餐店
    add_image(slide, img1_path, 609600, 1150000, 2700000, 1800000)

    # Label under image 1
    add_textbox(slide, 609600, 2970000, 2700000, 200000,
                "傳統模式：08:00–14:00 打烊", "374151", 13, False, PP_ALIGN.CENTER)

    # Image 2: AI自助24H吧台
    add_image(slide, img2_path, 609600, 3200000, 2700000, 1800000)

    # Label under image 2
    add_textbox(slide, 609600, 5010000, 2700000, 200000,
                "AI 自助：24小時不停收款", "00A896", 13, True, PP_ALIGN.CENTER)

    # 3. CENTER FINANCIAL TABLE
    center_left = 3500000

    # Big headline box
    hl = add_rect(slide, center_left, 1150000, 4600000, 400000, 'FFF8E8', 'C9A84C', 1)
    add_textbox(slide, center_left, 1150000, 4600000, 400000,
                "🏠 房租固定 → 多開就是純賺", "0D1B2A", 18, True, PP_ALIGN.CENTER)

    # Table background
    add_rect(slide, center_left, 1620000, 4600000, 3400000, 'F0F4F8')

    # Row 1 header
    add_rect(slide, center_left, 1620000, 4600000, 350000, '0D1B2A')
    add_textbox(slide, center_left, 1620000, 4600000, 350000,
                "每日收入試算（一間店）", "FFFFFF", 14, True, PP_ALIGN.CENTER)

    # BEFORE section header
    add_rect(slide, center_left, 1970000, 4600000, 260000, 'E8E8E8')
    add_textbox(slide, center_left + 50000, 1970000, 4600000 - 50000, 260000,
                "❌ 傳統 8H 模式", "374151", 12, False)

    # Row items - BEFORE
    add_two_col_row(slide, center_left, 2230000, 4600000, 190000,
                    "早餐 06-10 時", "NT$7,650", "374151", "374151", 11, 11)
    add_two_col_row(slide, center_left, 2420000, 4600000, 190000,
                    "早午餐 10-14 時", "NT$5,400", "374151", "374151", 11, 11)
    add_two_col_row(slide, center_left, 2610000, 4600000, 190000,
                    "14:00 後設備閒置", "NT$0", "374151", "EF4444", 11, 11)
    add_two_col_row(slide, center_left, 2800000, 4600000, 190000,
                    "傳統日收入合計", "NT$13,050", "0D1B2A", "0D1B2A", 12, 12, True, True)

    # Divider line
    add_rect(slide, center_left, 3050000, 4600000, 8000, 'C9A84C')

    # AFTER section header
    add_rect(slide, center_left, 3100000, 4600000, 250000, 'E0F5F2')
    add_textbox(slide, center_left + 50000, 3100000, 4600000 - 50000, 250000,
                "✅ AI 自助 24H 模式", "00C6AD", 12, True)

    # Row items - AFTER
    add_two_col_row(slide, center_left, 3350000, 4600000, 180000,
                    "早餐 06-10 時", "NT$7,650（同）", "374151", "374151", 11, 11)
    add_two_col_row(slide, center_left, 3530000, 4600000, 180000,
                    "早午餐 10-14 時", "NT$5,400（同）", "374151", "374151", 11, 11)
    add_two_col_row(slide, center_left, 3710000, 4600000, 180000,
                    "下午茶 14-18 時（自助咖啡）", "NT$3,900", "00A896", "00A896", 11, 11)
    add_two_col_row(slide, center_left, 3890000, 4600000, 180000,
                    "晚間 18-22 時（咖啡+輕食）", "NT$2,100", "00A896", "00A896", 11, 11)
    add_two_col_row(slide, center_left, 4070000, 4600000, 180000,
                    "AI 後日收入合計", "NT$19,050", "0D1B2A", "0D1B2A", 12, 12, True, True)

    # Cost section
    add_rect(slide, center_left, 4320000, 4600000, 220000, 'EEF2F7')
    add_textbox(slide, center_left + 40000, 4320000, 4600000 - 40000, 220000,
                "新增成本：電費+耗材 NT$700 ｜設備攤提 NT$440 ｜房租 NT$0 ｜人力 NT$0",
                "6B7280", 10, False)

    # Net gain highlight
    add_rect(slide, center_left, 4600000, 4600000, 400000, 'C9A84C')
    add_textbox(slide, center_left, 4600000, 4600000, 400000,
                "每日淨增盈：NT$5,860　月增：NT$17.6萬　年增：NT$213.9萬",
                "FFFFFF", 14, True, PP_ALIGN.CENTER)

    # 4. RIGHT Q&A COLUMN
    right_left = 8300000
    right_width = 3650000

    # Header
    add_rect(slide, right_left, 1150000, right_width, 350000, '0D1B2A')
    add_textbox(slide, right_left, 1150000, right_width, 350000,
                "📊 投資人問答", "FFFFFF", 16, True, PP_ALIGN.CENTER)

    # Q1
    add_rect(slide, right_left, 1550000, right_width, 260000, 'FFF0F0')
    add_textbox(slide, right_left + 40000, 1550000, right_width - 40000, 260000,
                "Q: 下午真的有人買嗎？數字是假設的？", "EF4444", 11, True)

    # A1
    add_rect(slide, right_left, 1810000, right_width, 300000, 'F0FFF4')
    add_textbox(slide, right_left + 40000, 1810000, right_width - 40000, 300000,
                "A: 麥味登實際驗證 — 導入後週訂單 +26%，非早餐時段佔比 32%。不是假設，是業界已驗證數據。",
                "16A34A", 10, False)

    # Q2
    add_rect(slide, right_left, 2160000, right_width, 260000, 'FFF0F0')
    add_textbox(slide, right_left + 40000, 2160000, right_width - 40000, 260000,
                "Q: 設備費 NT$80萬，回本要多久？", "EF4444", 11, True)

    # A2
    add_rect(slide, right_left, 2420000, right_width, 300000, 'F0FFF4')
    add_textbox(slide, right_left + 40000, 2420000, right_width - 40000, 300000,
                "A: 每日淨增 NT$5,860。NT$800,000 ÷ NT$5,860 = 137 天（約 4.5 個月）即回本。",
                "16A34A", 10, False)

    # Q3
    add_rect(slide, right_left, 2770000, right_width, 260000, 'FFF0F0')
    add_textbox(slide, right_left + 40000, 2770000, right_width - 40000, 260000,
                "Q: 沒員工值班，設備壞掉怎辦？", "EF4444", 11, True)

    # A3
    add_rect(slide, right_left, 3030000, right_width, 300000, 'F0FFF4')
    add_textbox(slide, right_left + 40000, 3030000, right_width - 40000, 300000,
                "A: 遠端 AI 監控即時告警，ACER 海柏特承諾 4 小時 SLA 到場維修。99.6% 可用率。",
                "16A34A", 10, False)

    # Divider
    add_rect(slide, right_left, 3380000, right_width, 6000, 'C9A84C')

    # Scale section header
    add_rect(slide, right_left, 3430000, right_width, 280000, 'EEF2F7')
    add_textbox(slide, right_left + 40000, 3430000, right_width - 40000, 280000,
                "規模效益（全面導入）", "0D1B2A", 13, True)

    # Scale item 1
    add_two_col_row(slide, right_left, 3750000, right_width, 280000,
                    "試點 50 間店", "NT$1.07億/年", "6B7280", "00C6AD", 11, 14, False, True)

    # Scale item 2
    add_two_col_row(slide, right_left, 4070000, right_width, 280000,
                    "全面 970 間", "NT$20.8億/年", "0D1B2A", "C9A84C", 11, 18, True, True)

    # Conclusion box
    add_rect(slide, right_left, 4420000, right_width, 500000, '0D1B2A')
    add_textbox_multiline(slide, right_left, 4420000, right_width, 500000, [
        ("回本僅需 4.5 個月", "FFFFFF", 20, True, PP_ALIGN.CENTER),
        ("之後每年每店多賺 NT$213 萬", "C9A84C", 13, False, PP_ALIGN.CENTER),
    ])

    print("  Slide 15 rebuilt successfully!")


def main():
    print("Loading presentation...")
    prs = Presentation(INPUT_PATH)
    print(f"Loaded {len(prs.slides)} slides")

    print("\nStep 1: Applying global bright theme to all slides...")
    process_all_slides(prs)

    print("\nStep 2: Rebuilding slide 15...")
    slide15 = prs.slides[14]  # 0-indexed
    rebuild_p15(slide15, IMG1_PATH, IMG2_PATH, IMG3_PATH)

    print(f"\nSaving to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print(f"\nDone! Output saved to:\n{OUTPUT_PATH}")


if __name__ == '__main__':
    main()
